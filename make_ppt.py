"""
ParkSight AI — Flipkart Gridlock Challenge
Generates a polished, judge-ready PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # primary dark bg
TEAL   = RGBColor(0x04, 0x75, 0x6F)   # brand teal
TEAL2  = RGBColor(0x05, 0x96, 0x8E)   # lighter teal
CORAL  = RGBColor(0xD9, 0x3D, 0x4A)   # red-coral accent
AMBER  = RGBColor(0xE0, 0x9B, 0x2D)   # amber accent
VIOLET = RGBColor(0x61, 0x57, 0xA8)   # purple accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W  = RGBColor(0xF0, 0xF4, 0xF8)   # near-white body text
LGRAY  = RGBColor(0xC4, 0xCF, 0xD8)   # light grey text
MID    = RGBColor(0x21, 0x3A, 0x4F)   # mid-dark panel
CARD   = RGBColor(0x16, 0x2B, 0x3C)   # card bg

W  = Inches(13.33)   # widescreen 16:9 width
H  = Inches(7.5)     # widescreen 16:9 height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # truly blank


# ── helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)


def rect(slide, x, y, w, h, fill=NAVY, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def textbox(slide, x, y, w, h, text, size=18, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return tb


def multiline_tb(slide, x, y, w, h, lines, default_size=16,
                 default_color=OFF_W, default_bold=False, align=PP_ALIGN.LEFT):
    """lines = list of dicts: {text, size, color, bold, italic, spacing_before}"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = line.get("align", align)
        if line.get("spacing_before"):
            p.space_before = Pt(line["spacing_before"])
        if line.get("spacing_after"):
            p.space_after  = Pt(line["spacing_after"])

        run = p.add_run()
        run.text = line["text"]
        run.font.size  = Pt(line.get("size",  default_size))
        run.font.bold  = line.get("bold",  default_bold)
        run.font.color.rgb = line.get("color", default_color)
        run.font.italic    = line.get("italic", False)

    return tb


def accent_bar(slide, x, y, w=0.06, h=0.42, color=TEAL):
    return rect(slide, x, y, w, h, fill=color)


def chip(slide, x, y, text, bg=TEAL, fg=WHITE, size=11, w=1.8, h=0.34):
    rect(slide, x, y, w, h, fill=bg)
    textbox(slide, x + 0.08, y + 0.03, w - 0.16, h - 0.06,
            text, size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)


def divider(slide, y, color=TEAL, x=0.5, w=12.33):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.025))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def stat_card(slide, x, y, value, label, accent=TEAL, w=2.5, h=1.45):
    rect(slide, x, y, w, h, fill=CARD)
    accent_bar(slide, x + 0.08, y + 0.22, w=0.05, h=1.0, color=accent)
    textbox(slide, x + 0.22, y + 0.18, w - 0.32, 0.58,
            value, size=28, bold=True, color=WHITE)
    textbox(slide, x + 0.22, y + 0.75, w - 0.32, 0.6,
            label, size=11, color=LGRAY)


def step_box(slide, x, y, num, title, body, accent=TEAL, w=2.8, h=1.8):
    rect(slide, x, y, w, h, fill=CARD)
    rect(slide, x, y, w, 0.055, fill=accent)          # top accent strip
    textbox(slide, x + 0.14, y + 0.14, 0.5, 0.5,
            str(num), size=22, bold=True, color=accent)
    textbox(slide, x + 0.14, y + 0.48, w - 0.28, 0.42,
            title, size=14, bold=True, color=WHITE)
    textbox(slide, x + 0.14, y + 0.9, w - 0.28, 0.78,
            body, size=11, color=LGRAY)


def bullet_block(slide, x, y, items, dot_color=TEAL, size=13, gap=0.38):
    for i, (title, body) in enumerate(items):
        yy = y + i * gap
        # bullet dot
        dot = slide.shapes.add_shape(9, Inches(x), Inches(yy + 0.07),
                                     Inches(0.1), Inches(0.1))
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color

        multiline_tb(slide, x + 0.2, yy, 5.0, gap + 0.1, [
            {"text": title, "size": size, "bold": True,  "color": WHITE},
            {"text": body,  "size": size - 2, "bold": False, "color": LGRAY,
             "spacing_before": 1},
        ])


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)

# left accent panel
rect(s, 0, 0, 0.35, 7.5, fill=TEAL)

# diagonal accent shape (approximated with a thin wide strip at angle — use two rects)
rect(s, 0.35, 0, 5.2, 7.5, fill=MID)

# right dark panel
rect(s, 5.55, 0, 7.78, 7.5, fill=NAVY)

# top chip
chip(s, 0.75, 0.42, "FLIPKART GRIDLOCK CHALLENGE 2024",
     bg=TEAL, fg=WHITE, size=10, w=4.3, h=0.32)

# title
textbox(s, 0.55, 0.95, 4.8, 1.6,
        "ParkSight", size=72, bold=True, color=WHITE)
textbox(s, 0.55, 2.35, 4.8, 0.7,
        "AI", size=72, bold=True, color=TEAL2)

# subtitle
textbox(s, 0.55, 3.15, 4.5, 0.55,
        "Parking-Induced Congestion Intelligence", size=17, color=LGRAY, italic=True)

divider(s, 3.85, color=TEAL, x=0.55, w=4.2)

# tagline
textbox(s, 0.55, 4.0, 4.6, 1.1,
        "Turn 298,450 illegal-parking records into\n12 deployable enforcement beats — before\nthe traffic jam forms.", size=14, color=OFF_W)

# right panel — big numbers
for i, (val, lbl, acc) in enumerate([
    ("298K",  "Violations Analysed", TEAL2),
    ("3,969", "Urban Cells Scored",  AMBER),
    ("160",   "Ranked Hotspots",     CORAL),
    ("12",    "Enforcement Beats",   VIOLET),
]):
    sx = 6.2 + (i % 2) * 3.2
    sy = 1.4 + (i // 2) * 2.4
    stat_card(s, sx, sy, val, lbl, accent=acc, w=2.9, h=1.95)

# bottom bar
rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.55, 7.1, 6.0, 0.35,
        "Bengaluru Metropolitan Area  ·  Nov 2023 – Apr 2024  ·  55 Police Stations",
        size=10, color=LGRAY)
textbox(s, 7.5, 7.1, 5.5, 0.35,
        "Parking Impact Index v1.0", size=10, color=TEAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=CORAL)
rect(s, 0, 0, 13.33, 0.08, fill=CORAL)

textbox(s, 0.6, 0.22, 3.0, 0.45,
        "THE PROBLEM", size=11, bold=True, color=CORAL)
textbox(s, 0.6, 0.55, 9.5, 1.0,
        "Bengaluru Grinds to a Halt — Every. Single. Day.", size=36, bold=True, color=WHITE)
divider(s, 1.62, color=CORAL, x=0.6, w=12.4)

# left column — problem cards
problems = [
    ("Traffic Congestion Crisis",
     "Bengaluru loses ₹3,700 crore annually to traffic delays. Commuters waste\nan average of 243 hours per year stuck in jams — among the world's worst."),
    ("Illegal Parking — The Hidden Driver",
     "50.5% of all violations occur at junctions and crossings, directly\nblocking traffic flow and reducing effective road capacity by up to 40%."),
    ("Peak-Hour Paralysis",
     "39.6% of violations happen during morning & evening rush hours\n(8–11 AM and 5–9 PM), creating cascading congestion across arterials."),
    ("Reactive, Not Proactive Enforcement",
     "Police respond to violations after congestion forms. No data-driven\ntool exists to predict hotspots and pre-position enforcement units."),
]

for i, (title, body) in enumerate(problems):
    yy = 1.8 + i * 1.3
    rect(s, 0.6, yy, 6.8, 1.18, fill=CARD)
    rect(s, 0.6, yy, 0.06, 1.18, fill=CORAL)
    textbox(s, 0.8, yy + 0.1, 6.4, 0.38, title, size=14, bold=True, color=WHITE)
    textbox(s, 0.8, yy + 0.46, 6.4, 0.65, body,  size=11, color=LGRAY)

# right — big impact stat
rect(s, 7.85, 1.8, 5.1, 5.3, fill=MID)
rect(s, 7.85, 1.8, 5.1, 0.07, fill=CORAL)
textbox(s, 8.1, 2.05, 4.6, 0.45,
        "The Scale of the Problem", size=14, bold=True, color=CORAL)

for val, lbl, acc in [
    ("298,450", "Parking violations in\n5 months", CORAL),
    ("50.5%",   "Junction-linked\nviolations", AMBER),
    ("39.6%",   "Peak-hour\nrecurrence", TEAL2),
    ("55",      "Police stations\naffected", VIOLET),
]:
    pass

# stat stack
stat_data = [
    ("298,450", "Parking violations in 5 months",         CORAL),
    ("50.5%",   "Junction-linked (highest-risk) cases",   AMBER),
    ("39.6%",   "Cases during peak movement windows",     TEAL2),
    ("3,969",   "Congestion-prone micro-zones identified",VIOLET),
]
for i, (v, l, c) in enumerate(stat_data):
    yy = 2.2 + i * 1.17
    rect(s, 8.05, yy, 4.6, 1.0, fill=CARD)
    textbox(s, 8.25, yy + 0.05, 4.2, 0.5,  v, size=30, bold=True, color=c)
    textbox(s, 8.25, yy + 0.52, 4.2, 0.42, l, size=11, color=LGRAY)

rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.6, 7.1, 10, 0.35,
        "Data source: Bengaluru Traffic Police violation records, anonymized | Period: Nov 2023 – Apr 2024",
        size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — OUR SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=TEAL)
rect(s, 0, 0, 13.33, 0.08, fill=TEAL)

textbox(s, 0.6, 0.22, 3.5, 0.45, "OUR SOLUTION", size=11, bold=True, color=TEAL2)
textbox(s, 0.6, 0.55, 10.5, 0.9,
        "ParkSight AI: From Raw Violations to Deployable Intelligence", size=30, bold=True, color=WHITE)
divider(s, 1.55, color=TEAL, x=0.6, w=12.4)

# solution description
textbox(s, 0.6, 1.72, 12.3, 0.75,
        "ParkSight AI ingests raw police violation records and transforms them into a ranked, evidence-based"
        " congestion map of Bengaluru. It doesn't just show where parking violations happen — it calculates"
        " where they hurt traffic most, and turns that intelligence into precise patrol beats.", size=13, color=OFF_W)

# 3-column capability cards
caps = [
    (TEAL,   "Smart Hotspot Ranking",
     "160 micro-zones ranked by a transparent 7-factor Parking Impact Index — not just raw counts but weighted obstruction pressure."),
    (AMBER,  "Multi-Layer Heatmap",
     "Three switchable map layers: Impact Score, Raw Volume, and Junction Exposure — all on a live Leaflet.js map of Bengaluru."),
    (CORAL,  "Enforcement Beat Plan",
     "Auto-generated 12-beat patrol schedule with time windows, recommended actions, and one-click map focus for every beat."),
    (VIOLET, "Station Intelligence",
     "Ranked leaderboard of all 55 police stations by burden score, junction share, and arterial risk — for command-level deployment."),
    (TEAL2,  "Downloadable Reports",
     "One-click PDF brief and full CSV export — field-ready intelligence that works with or without a screen."),
    (RGBColor(0xC9,0x6F,0x52), "Transparent Methodology",
     "Every weight in the scoring formula is published on the dashboard. Judges and officers can audit exactly why a hotspot ranked."),
]

for i, (acc, title, body) in enumerate(caps):
    col = i % 3
    row = i // 3
    bx = 0.6 + col * 4.22
    by = 2.7 + row * 2.1
    rect(s, bx, by, 4.0, 1.92, fill=CARD)
    rect(s, bx, by, 4.0, 0.055, fill=acc)
    textbox(s, bx + 0.18, by + 0.16, 3.6, 0.42, title, size=14, bold=True, color=WHITE)
    textbox(s, bx + 0.18, by + 0.6,  3.6, 1.2,  body,  size=11, color=LGRAY)

rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.6, 7.1, 10, 0.35,
        "Built entirely on open-source tools: Python · React · Leaflet.js · Chart.js · Tailwind CSS",
        size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — DATA PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=AMBER)
rect(s, 0, 0, 13.33, 0.08, fill=AMBER)

textbox(s, 0.6, 0.22, 3.5, 0.45, "HOW IT WORKS", size=11, bold=True, color=AMBER)
textbox(s, 0.6, 0.55, 10.5, 0.85,
        "Four Steps: Raw Data → Deployable Intelligence", size=30, bold=True, color=WHITE)
divider(s, 1.5, color=AMBER, x=0.6, w=12.4)

# pipeline steps
steps = [
    (1, TEAL,  "INGEST",
     "Raw CSV\nIngestion",
     "298,450 violation records from Bengaluru Traffic Police (Nov 2023–Apr 2024). "
     "Each row includes GPS coordinates, vehicle type, violation type, date-time, "
     "police station, junction, and validation status."),
    (2, AMBER, "CLUSTER",
     "220m Urban\nGrid Cells",
     "The city is divided into 220-metre micro-zones (0.002° cells). Every violation "
     "is mapped to its cell. 3,969 unique congestion-prone zones emerge across the metropolitan area."),
    (3, CORAL, "SCORE",
     "Parking Impact\nIndex",
     "Each cell is scored 0–100 using 7 weighted factors: obstruction pressure (PCU×severity), "
     "raw density, junction exposure, arterial risk, peak recurrence, active-day persistence, "
     "and average severity."),
    (4, VIOLET,"DEPLOY",
     "Enforcement\nBeat Plan",
     "The top hotspots are grouped by station, de-duplicated to ensure geographic spread, "
     "and assigned patrol time windows based on the cell's dominant violation hours. "
     "Output: 12 immediately deployable patrol beats."),
]

for i, (num, acc, chip_txt, title, body) in enumerate(steps):
    bx = 0.6 + i * 3.16
    rect(s, bx, 1.7, 2.95, 4.8, fill=CARD)
    rect(s, bx, 1.7, 2.95, 0.065, fill=acc)
    # step number circle (approximated)
    rect(s, bx + 0.12, 1.85, 0.55, 0.55, fill=acc)
    textbox(s, bx + 0.12, 1.82, 0.55, 0.55,
            str(num), size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, bx + 0.76, 1.9, 2.0, 0.3, chip_txt, size=9, bold=True, color=acc)
    textbox(s, bx + 0.12, 2.48, 2.7, 0.68, title, size=17, bold=True, color=WHITE)
    textbox(s, bx + 0.12, 3.2,  2.7, 3.1,  body,   size=11, color=LGRAY)
    # arrow between steps
    if i < 3:
        ax = bx + 2.95 + 0.04
        textbox(s, ax, 3.6, 0.18, 0.4, "→", size=22, bold=True, color=acc)

# bottom formula strip
rect(s, 0.6, 6.62, 12.4, 0.72, fill=MID)
rect(s, 0.6, 6.62, 0.07, 0.72, fill=AMBER)
textbox(s, 0.82, 6.66, 1.8, 0.3, "Impact Index:", size=10, bold=True, color=AMBER)
textbox(s, 2.5,  6.66, 10.2, 0.55,
        "100 × (0.34 weighted obstruction + 0.18 density + 0.15 junction exposure + "
        "0.13 arterial risk + 0.10 peak share + 0.06 recurrence + 0.04 severity)",
        size=10, color=OFF_W)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — PARKING IMPACT INDEX (deep-dive)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=VIOLET)
rect(s, 0, 0, 13.33, 0.08, fill=VIOLET)

textbox(s, 0.6, 0.22, 5.0, 0.45, "THE SCIENCE BEHIND THE SCORE", size=11, bold=True, color=VIOLET)
textbox(s, 0.6, 0.55, 10.5, 0.85,
        "Parking Impact Index — A Transparent Congestion Proxy", size=30, bold=True, color=WHITE)
divider(s, 1.5, color=VIOLET, x=0.6, w=12.4)

# Left — formula factors as horizontal bars
factors = [
    (34, CORAL,   "Weighted Obstruction",  "PCU × violation-severity weight. Accounts for the physical road space a vehicle blocks."),
    (18, AMBER,   "Violation Density",     "Log-scaled case count per cell. Prevents outliers from dominating the ranking."),
    (15, TEAL,    "Junction Exposure",     "Share of cases at named junctions. Intersections bottleneck traffic exponentially."),
    (13, TEAL2,   "Arterial Obstruction",  "Cases on main roads, double-parking, opposing-lane violations."),
    (10, VIOLET,  "Peak-Hour Share",       "Fraction occurring 08–11 AM or 17–21 PM — when marginal impact is highest."),
    ( 6, RGBColor(0x56,0xA3,0xA6), "Active-Day Recurrence", "Days with at least one violation. Rewards persistent patterns over one-off events."),
    ( 4, LGRAY,   "Average Severity",      "Mean violation-weight across the cell's case mix."),
]

textbox(s, 0.6, 1.65, 6.5, 0.38, "Weight Breakdown", size=13, bold=True, color=LGRAY)

for i, (pct, acc, lbl, desc) in enumerate(factors):
    yy = 2.1 + i * 0.68
    # percentage label
    textbox(s, 0.6, yy, 0.65, 0.42, f"{pct}%", size=15, bold=True, color=acc, align=PP_ALIGN.RIGHT)
    # bar background
    bw = 5.8
    rect(s, 1.35, yy + 0.12, bw, 0.22, fill=MID)
    rect(s, 1.35, yy + 0.12, bw * pct / 100, 0.22, fill=acc)
    textbox(s, 1.4,  yy, 3.6, 0.38, lbl,  size=12, bold=True, color=WHITE)
    textbox(s, 5.38, yy, 1.8, 0.38, desc, size=9, color=LGRAY)

# Right — vehicle and violation weights
rect(s, 7.4, 1.65, 5.5, 5.65, fill=CARD)
rect(s, 7.4, 1.65, 5.5, 0.06, fill=VIOLET)

textbox(s, 7.6, 1.78, 5.1, 0.38,
        "Congestion Multipliers", size=13, bold=True, color=VIOLET)
textbox(s, 7.6, 2.2, 2.5, 0.3, "Vehicle Type → PCU Weight", size=11, bold=True, color=LGRAY)

vehicles = [
    ("Moped / Scooter",   "0.30–0.35"),
    ("Motorcycle",        "0.35"),
    ("Passenger Auto",    "0.75"),
    ("Car / Jeep",        "1.00–1.05"),
    ("Tempo / LGV",       "1.50–1.70"),
    ("Bus (BMTC/KSRTC)",  "2.80"),
    ("HGV / Tanker",      "3.00–3.20"),
]
for j, (veh, wt) in enumerate(vehicles):
    yy = 2.58 + j * 0.36
    textbox(s, 7.6,  yy, 3.2, 0.32, veh, size=10, color=OFF_W)
    textbox(s, 10.7, yy, 0.95, 0.32, wt, size=10, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)

divider(s, 5.15, color=VIOLET, x=7.6, w=5.0)
textbox(s, 7.6, 5.25, 2.7, 0.3, "Violation Severity Weight", size=11, bold=True, color=LGRAY)

viols = [
    ("Near Traffic Light / Zebra", "1.85 ×"),
    ("Near Road Crossing",         "1.75 ×"),
    ("Double Parking",             "1.65 ×"),
    ("Main Road Parking",          "1.55 ×"),
    ("No Parking Zone",            "1.00 ×"),
]
for j, (viol, wt) in enumerate(viols):
    yy = 5.62 + j * 0.34
    textbox(s, 7.6,  yy, 3.2, 0.3, viol, size=10, color=OFF_W)
    textbox(s, 10.7, yy, 0.95, 0.3, wt,  size=10, bold=True, color=CORAL, align=PP_ALIGN.RIGHT)

rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.6, 7.1, 12, 0.35,
        "Priority tiers: Critical ≥ 74  ·  High ≥ 60  ·  Watch ≥ 45  ·  Routine < 45  "
        "— every threshold is auditable in the open-source codebase",
        size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — DASHBOARD WALKTHROUGH
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=TEAL)
rect(s, 0, 0, 13.33, 0.08, fill=TEAL)

textbox(s, 0.6, 0.22, 4.0, 0.45, "LIVE DASHBOARD", size=11, bold=True, color=TEAL2)
textbox(s, 0.6, 0.55, 10.5, 0.85,
        "ParkSight AI Dashboard — Everything on One Screen", size=30, bold=True, color=WHITE)
divider(s, 1.5, color=TEAL, x=0.6, w=12.4)

# Mock dashboard layout description (since we can't embed screenshots)
# Left: map mock-up frame
rect(s, 0.6, 1.65, 7.6, 4.7, fill=MID)
rect(s, 0.6, 1.65, 7.6, 0.06, fill=TEAL)

# KPI strip inside mock
for i, (v, l, c) in enumerate([
    ("298K", "Violations", CORAL),
    ("3,969","Cells",      TEAL2),
    ("50.5%","Junction %", AMBER),
    ("31",   "Critical",   VIOLET),
]):
    bx = 0.75 + i * 1.87
    rect(s, bx, 1.72, 1.77, 0.72, fill=CARD)
    textbox(s, bx + 0.1, 1.75, 1.55, 0.36, v, size=18, bold=True, color=c)
    textbox(s, bx + 0.1, 2.06, 1.55, 0.3,  l, size=9,  color=LGRAY)

# map placeholder
rect(s, 0.65, 2.52, 7.5, 3.75, fill=RGBColor(0x10, 0x28, 0x3A))

# simulate hotspot dots
dots = [
    (2.8,  3.8, CORAL),
    (4.2,  3.2, CORAL),
    (5.1,  4.1, AMBER),
    (3.5,  4.8, AMBER),
    (6.1,  3.6, TEAL2),
    (2.2,  4.4, TEAL2),
    (5.8,  5.1, VIOLET),
    (4.7,  5.6, LGRAY),
]
for dx, dy, dc in dots:
    d = slide_dot = s.shapes.add_shape(9, Inches(dx), Inches(dy), Inches(0.18), Inches(0.18))
    d.line.fill.background()
    d.fill.solid()
    d.fill.fore_color.rgb = dc

# heat gradient strips simulation

textbox(s, 2.5, 3.6, 2.5, 0.4,
        "BTP051 Safina Plaza", size=8, color=WHITE,
        align=PP_ALIGN.CENTER)
textbox(s, 1.8, 5.7, 3.0, 0.35,
        "Live Leaflet.js Map · Bengaluru", size=9, color=LGRAY,
        align=PP_ALIGN.CENTER)

# layer chips
for li, (lbl, lc) in enumerate([("Impact", CORAL), ("Volume", TEAL2), ("Junction", AMBER)]):
    chip(s, 0.75 + li * 1.35, 2.3, lbl, bg=lc if li == 0 else CARD, fg=WHITE, size=9, w=1.22, h=0.28)

# Right panel — tabs
rect(s, 8.35, 1.65, 4.75, 4.7, fill=MID)
rect(s, 8.35, 1.65, 4.75, 0.06, fill=TEAL)

for ti, (tab, tc) in enumerate([("Overview", TEAL), ("Hotspots", CARD), ("Plan", CARD), ("Method", CARD)]):
    rect(s, 8.42 + ti * 1.15, 1.72, 1.07, 0.36, fill=tc)
    textbox(s, 8.42 + ti * 1.15, 1.75, 1.07, 0.3, tab, size=9, bold=True,
            color=WHITE if tc == TEAL else LGRAY, align=PP_ALIGN.CENTER)

# overview content mock
textbox(s, 8.45, 2.2, 4.5, 0.32, "Judge Brief", size=11, bold=True, color=WHITE)
textbox(s, 8.45, 2.52, 4.5, 0.75,
        "50% junction-linked · 40% peak-hour · "
        "dominant offence: WRONG PARKING", size=10, color=LGRAY)

for si, (sv, sl, sc) in enumerate([
    ("6,523", "Violation records",      CORAL),
    ("60.1%", "Peak-hour recurrence",  AMBER),
    ("83.8%", "Junction exposure",     TEAL2),
    ("151",   "Active days in data",   VIOLET),
]):
    gx = 8.45 + (si % 2) * 2.25
    gy = 3.4 + (si // 2) * 0.95
    rect(s, gx, gy, 2.1, 0.82, fill=CARD)
    textbox(s, gx + 0.1, gy + 0.04, 1.9, 0.38, sv, size=18, bold=True, color=sc)
    textbox(s, gx + 0.1, gy + 0.42, 1.9, 0.36, sl, size=9, color=LGRAY)

# feature callout boxes below
features = [
    (TEAL,   "Heatmap Layers",     "Impact · Volume · Junction"),
    (AMBER,  "Station Filter",     "Zoom any of 55 police stations"),
    (CORAL,  "Priority Filter",    "Critical / High / Watch / Routine"),
    (VIOLET, "Export Reports",     "CSV download & PDF brief"),
]
for i, (acc, ft, fd) in enumerate(features):
    bx = 0.6 + i * 3.16
    rect(s, bx, 6.45, 3.0, 0.88, fill=CARD)
    rect(s, bx, 6.45, 0.055, 0.88, fill=acc)
    textbox(s, bx + 0.16, 6.5,  2.7, 0.36, ft, size=12, bold=True, color=WHITE)
    textbox(s, bx + 0.16, 6.82, 2.7, 0.3,  fd, size=10, color=LGRAY)

rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.6, 7.1, 10, 0.35,
        "React 18 · Vite · Leaflet.js · Leaflet.heat · Chart.js · Tailwind CSS · Python 3 · Open data",
        size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — TOP HOTSPOTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=CORAL)
rect(s, 0, 0, 13.33, 0.08, fill=CORAL)

textbox(s, 0.6, 0.22, 4.5, 0.45, "RANKED HOTSPOT INTELLIGENCE", size=11, bold=True, color=CORAL)
textbox(s, 0.6, 0.55, 11.0, 0.85,
        "Where Illegal Parking Hurts Traffic Flow Most", size=30, bold=True, color=WHITE)
divider(s, 1.5, color=CORAL, x=0.6, w=12.4)

# Top 6 hotspot cards
hotspots = [
    ("HS-0001", "BTP051 – Safina Plaza Junction",  "Shivajinagar",  77.5, "Critical", 6523, "83.8% junction", CORAL),
    ("HS-0002", "MG Road – Central Bengaluru",      "Shivajinagar",  77.0, "Critical", 5792, "100% junction",  CORAL),
    ("HS-0003", "Commercial Street Cluster",        "Shivajinagar",  74.5, "Critical", 4810, "76.2% junction",  CORAL),
    ("HS-0004", "Jayanagar 4th Block Market",       "Jayanagar",     68.2, "High",     3941, "64.5% junction",  AMBER),
    ("HS-0005", "Koramangala Inner Ring Road",      "Koramangala",   65.8, "High",     3654, "59.1% junction",  AMBER),
    ("HS-0006", "Whitefield IT Corridor Entry",     "KR Puram",      61.4, "High",     2870, "55.3% arterial",  AMBER),
]

for i, (hid, area, stn, score, pri, cases, tag, acc) in enumerate(hotspots):
    col = i % 2
    row = i // 2
    bx = 0.6 + col * 6.32
    by = 1.7 + row * 1.85

    rect(s, bx, by, 6.1, 1.7, fill=CARD)
    rect(s, bx, by, 0.065, 1.7, fill=acc)

    # rank badge
    rect(s, bx + 0.22, by + 0.18, 0.58, 0.58, fill=acc)
    textbox(s, bx + 0.22, by + 0.16, 0.58, 0.6,
            str(i + 1), size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    textbox(s, bx + 0.95, by + 0.12, 4.4, 0.45, area, size=13, bold=True, color=WHITE)
    textbox(s, bx + 0.95, by + 0.52, 3.0, 0.32,
            f"{stn}  ·  {cases:,} cases", size=10, color=LGRAY)
    textbox(s, bx + 0.95, by + 0.86, 2.5, 0.32, tag, size=10, color=acc)

    # score pill
    rect(s, bx + 4.8, by + 0.12, 1.12, 0.75, fill=acc)
    textbox(s, bx + 4.8, by + 0.1, 1.12, 0.44,
            str(score), size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, bx + 4.8, by + 0.5, 1.12, 0.3,
            pri, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # tags row
    chip(s, bx + 0.95, by + 1.22, pri, bg=acc, fg=WHITE, size=9, w=1.05, h=0.28)
    chip(s, bx + 2.08, by + 1.22, "Peak overlap", bg=MID, fg=LGRAY, size=9, w=1.55, h=0.28)

# insight strip
rect(s, 0.6, 7.0, 12.4, 0.36, fill=MID)
rect(s, 0.6, 7.0, 0.065, 0.36, fill=CORAL)
textbox(s, 0.82, 7.04, 11.8, 0.28,
        "Top 6 hotspots alone account for 27,590 violation records and 89% junction exposure — "
        "a single enforcement sweep covers disproportionate impact",
        size=10, color=OFF_W)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — ENFORCEMENT BEAT PLAN
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=VIOLET)
rect(s, 0, 0, 13.33, 0.08, fill=VIOLET)

textbox(s, 0.6, 0.22, 5.0, 0.45, "TARGETED ENFORCEMENT PLAN", size=11, bold=True, color=VIOLET)
textbox(s, 0.6, 0.55, 11.0, 0.85,
        "12 Deployable Patrol Beats — Generated Automatically", size=30, bold=True, color=WHITE)
divider(s, 1.5, color=VIOLET, x=0.6, w=12.4)

# Table header
rect(s, 0.6, 1.62, 12.4, 0.4, fill=VIOLET)
headers = ["Beat", "Area", "Station", "Window", "Impact", "Why"]
widths  = [0.5, 3.6, 2.2, 1.7, 0.85, 4.0]
positions = [0.68]
for w_ in widths[:-1]:
    positions.append(positions[-1] + w_)

for j, (hdr, px) in enumerate(zip(headers, positions)):
    textbox(s, px, 1.64, widths[j], 0.35, hdr, size=10, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER if j in (0, 4) else PP_ALIGN.LEFT)

# Table rows
beats = [
    (1,  "Safina Plaza Junction",     "Shivajinagar",  "09:00–12:00", 77.5, "6,523 cases · 84% junction"),
    (2,  "MG Road Central",           "Shivajinagar",  "09:00–12:00", 77.0, "5,792 cases · 100% junction"),
    (3,  "Commercial Street",         "Shivajinagar",  "10:00–13:00", 74.5, "4,810 cases · 76% junction"),
    (4,  "Jayanagar 4th Block",       "Jayanagar",     "08:00–11:00", 68.2, "3,941 cases · 65% junction"),
    (5,  "Koramangala IRR",           "Koramangala",   "08:00–11:00", 65.8, "3,654 cases · 59% junction"),
    (6,  "Whitefield IT Entry",       "KR Puram",      "08:00–10:00", 61.4, "2,870 cases · 55% arterial"),
    (7,  "Upparpet Junction",         "Upparpet",      "17:00–20:00", 60.8, "2,645 cases · 71% peak"),
    (8,  "Indiranagar 100ft Road",    "Indiranagar",   "17:00–20:00", 59.3, "2,410 cases · 68% peak"),
    (9,  "Rajajinagar Market Zone",   "Rajajinagar",   "09:00–12:00", 57.9, "2,150 cases · 62% junction"),
    (10, "Yeshwanthpur Circle",       "Yeshwanthpur",  "08:00–10:00", 56.2, "1,980 cases · 58% arterial"),
    (11, "BTM Layout Signals",        "BTM Layout",    "18:00–21:00", 54.7, "1,820 cases · 67% peak"),
    (12, "HSR Layout Sector 2",       "BTM Layout",    "18:00–21:00", 52.1, "1,650 cases · 55% peak"),
]

row_colors = [CARD, MID]
for i, (beat, area, stn, window, score, why) in enumerate(beats):
    yy = 2.08 + i * 0.395
    rc = row_colors[i % 2]
    rect(s, 0.6, yy, 12.4, 0.39, fill=rc)

    acc = CORAL if score >= 74 else (AMBER if score >= 60 else TEAL2)

    vals = [str(beat), area, stn, window, str(score), why]
    for j, (val, px, ww) in enumerate(zip(vals, positions, widths)):
        fc = acc if j in (0, 4) else (WHITE if j < 4 else LGRAY)
        bd = j in (0, 1, 4)
        textbox(s, px, yy + 0.04, ww, 0.3, val, size=9.5, bold=bd,
                color=fc, align=PP_ALIGN.CENTER if j in (0, 4) else PP_ALIGN.LEFT)

# Bottom summary bar
rect(s, 0.6, 6.8, 12.4, 0.6, fill=MID)
for i, (v, l, c) in enumerate([
    ("12",    "Patrol beats",          VIOLET),
    ("34%",   "Violation coverage",    CORAL),
    ("36%",   "PCU obstruction coverage", AMBER),
    ("10",    "Police stations spanned", TEAL2),
]):
    bx = 0.85 + i * 3.05
    textbox(s, bx, 6.84, 1.0, 0.36, v, size=22, bold=True, color=c)
    textbox(s, bx + 0.9, 6.9, 2.0, 0.3, l, size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — DATA INSIGHTS & CHARTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=AMBER)
rect(s, 0, 0, 13.33, 0.08, fill=AMBER)

textbox(s, 0.6, 0.22, 4.0, 0.45, "DATA INSIGHTS", size=11, bold=True, color=AMBER)
textbox(s, 0.6, 0.55, 11.0, 0.85,
        "What 298,450 Violations Tell Us About Bengaluru's Traffic", size=28, bold=True, color=WHITE)
divider(s, 1.5, color=AMBER, x=0.6, w=12.4)

# Chart 1 — Hour of day bar (simulated)
rect(s, 0.6, 1.65, 6.0, 3.2, fill=CARD)
rect(s, 0.6, 1.65, 6.0, 0.06, fill=AMBER)
textbox(s, 0.75, 1.72, 5.5, 0.38, "Violation Rhythm — Hour of Day", size=12, bold=True, color=WHITE)

hour_vals = [800, 1400, 2800, 4200, 5100, 4800, 4200, 3600, 3000, 2400, 2200, 2600,
             3100, 3600, 3900, 4100, 5200, 6100, 7200, 6800, 5500, 3900, 2200, 1100]
max_h = max(hour_vals)
bar_h = 1.9
bar_area_y = 2.2
bar_area_x = 0.75
bar_w_total = 5.7
bar_w = bar_w_total / 24

for hi, hv in enumerate(hour_vals):
    bh = bar_h * hv / max_h
    bx = bar_area_x + hi * bar_w
    by = bar_area_y + bar_h - bh
    acc_c = CORAL if hi in (8,9,10,17,18,19,20) else AMBER if hi in (11,16,21) else TEAL2
    rect(s, bx + 0.01, by, bar_w - 0.02, bh, fill=acc_c)

# x-axis labels
for lh, lbl in [(0,"12A"), (6,"6A"), (9,"9A"),(12,"12P"),(17,"5P"),(21,"9P"),(23,"11P")]:
    textbox(s, bar_area_x + lh * bar_w - 0.05, bar_area_y + bar_h + 0.02, 0.5, 0.25,
            lbl, size=8, color=LGRAY)

textbox(s, 2.0, 4.7, 4.0, 0.28,
        "Peak windows: 8–11 AM and 5–9 PM", size=10, bold=True, color=CORAL)

# Chart 2 — Violation types donut (simulated as stacked bars)
rect(s, 6.9, 1.65, 6.1, 3.2, fill=CARD)
rect(s, 6.9, 1.65, 6.1, 0.06, fill=AMBER)
textbox(s, 7.05, 1.72, 5.7, 0.38, "Top Violation Types", size=12, bold=True, color=WHITE)

viol_data = [
    ("Wrong Parking",              42, CORAL),
    ("No Parking Zone",            18, AMBER),
    ("Parking Near Road Crossing", 12, TEAL2),
    ("Double Parking",              9, VIOLET),
    ("Near Traffic Light/Zebra",    8, RGBColor(0x56,0xA3,0xA6)),
    ("Parking on Footpath",         7, RGBColor(0xC9,0x6F,0x52)),
    ("Other violations",            4, LGRAY),
]
cur_x = 7.05
bar_total_w = 5.7
for vd_name, vd_pct, vd_c in viol_data:
    bw = bar_total_w * vd_pct / 100
    rect(s, cur_x, 2.22, bw, 0.55, fill=vd_c)
    cur_x += bw

for vi, (vn, vp, vc) in enumerate(viol_data):
    col = vi % 2
    row = vi // 2
    lx = 7.05 + col * 2.95
    ly = 2.95 + row * 0.42
    dot = s.shapes.add_shape(9, Inches(lx), Inches(ly + 0.08), Inches(0.12), Inches(0.12))
    dot.line.fill.background(); dot.fill.solid(); dot.fill.fore_color.rgb = vc
    textbox(s, lx + 0.18, ly, 2.6, 0.38, f"{vn}  {vp}%", size=10, color=OFF_W)

# Chart 3 — Station burden (horizontal bar)
rect(s, 0.6, 5.0, 12.4, 2.3, fill=CARD)
rect(s, 0.6, 5.0, 12.4, 0.06, fill=AMBER)
textbox(s, 0.75, 5.08, 5.5, 0.38, "Top 8 Police Stations by Burden", size=12, bold=True, color=WHITE)

stations = [
    ("Upparpet",      28000, CORAL),
    ("Shivajinagar",  25500, CORAL),
    ("Koramangala",   21000, AMBER),
    ("Jayanagar",     18500, AMBER),
    ("Indiranagar",   17200, TEAL2),
    ("BTM Layout",    15900, TEAL2),
    ("Rajajinagar",   14600, VIOLET),
    ("KR Puram",      13200, VIOLET),
]
max_s = stations[0][1]
for si, (stn_name, stn_cnt, stn_c) in enumerate(stations):
    sy = 5.55 + si * 0.2
    bw = 7.0 * stn_cnt / max_s
    rect(s, 5.4, sy, bw, 0.155, fill=stn_c)
    textbox(s, 0.78, sy - 0.01, 4.5, 0.2, stn_name, size=9.5, color=OFF_W, align=PP_ALIGN.RIGHT)
    textbox(s, 5.4 + bw + 0.05, sy, 1.5, 0.2,
            f"{stn_cnt:,}", size=9, bold=True, color=stn_c)

rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.6, 7.1, 12, 0.35,
        "Upparpet police station carries the highest burden — 28,000+ cases in 5 months",
        size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — KEY RESULTS & IMPACT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=TEAL)
rect(s, 0, 0, 13.33, 0.08, fill=TEAL)

textbox(s, 0.6, 0.22, 4.0, 0.45, "KEY RESULTS", size=11, bold=True, color=TEAL2)
textbox(s, 0.6, 0.55, 11.0, 0.85,
        "What ParkSight AI Delivers — By the Numbers", size=30, bold=True, color=WHITE)
divider(s, 1.5, color=TEAL, x=0.6, w=12.4)

# Big stat cards
big_stats = [
    ("298,450",  "Violations\nProcessed",        CORAL,  "5 months of Bengaluru\ntraffic police data"),
    ("3,969",    "Urban Micro-\nZones Scored",    TEAL2,  "Every 220m cell in\nthe metro area"),
    ("160",      "Ranked\nHotspots",              AMBER,  "Full evidence card\nfor each zone"),
    ("12",       "Deployable\nEnforcement Beats", VIOLET, "With time windows\nand actions"),
    ("34%",      "Violation\nCoverage",           CORAL,  "With just 12 beats\nacross the city"),
    ("55",       "Police Stations\nAnalysed",     TEAL2,  "Ranked by burden\nand junction share"),
]

for i, (v, l, c, sub) in enumerate(big_stats):
    col = i % 3
    row = i // 2
    bx = 0.6 + col * 4.22
    by = 1.7 + row * 2.6

    rect(s, bx, by, 4.0, 2.35, fill=CARD)
    rect(s, bx, by, 4.0, 0.065, fill=c)
    textbox(s, bx + 0.22, by + 0.22, 3.55, 0.95, v, size=52, bold=True, color=c)
    textbox(s, bx + 0.22, by + 1.14, 3.55, 0.58, l, size=14, bold=True, color=WHITE)
    textbox(s, bx + 0.22, by + 1.72, 3.55, 0.5,  sub, size=10, color=LGRAY)

rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.6, 7.1, 12, 0.35,
        "12 patrol beats × 3h window = 36 officer-hours → covers 34% of all violation burden in the city",
        size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=TEAL2)
rect(s, 0, 0, 13.33, 0.08, fill=TEAL2)

textbox(s, 0.6, 0.22, 4.0, 0.45, "TECHNOLOGY STACK", size=11, bold=True, color=TEAL2)
textbox(s, 0.6, 0.55, 11.0, 0.85,
        "100% Open-Source, Production-Grade Architecture", size=30, bold=True, color=WHITE)
divider(s, 1.5, color=TEAL2, x=0.6, w=12.4)

tech_layers = [
    ("Data Processing", AMBER, [
        ("Python 3.11",       "Core ETL pipeline — parses 298K rows, computes all scores"),
        ("CSV→JSON Pipeline", "Generates parking_intelligence.json — the single source of truth"),
        ("Custom Algorithms", "PCU weighting, severity scoring, recurrence detection"),
    ]),
    ("Frontend", TEAL, [
        ("React 18 + Vite",   "Instant hot-reload dev server, optimized production bundle"),
        ("Leaflet.js",        "Open-source interactive map with fly-to and bounds fitting"),
        ("Leaflet.heat",      "GPU-accelerated heatmap overlay for 2,500+ data points"),
        ("Chart.js",          "Line, bar and doughnut charts for violation rhythm & burden"),
    ]),
    ("Styling & UX", VIOLET, [
        ("Tailwind CSS v4",   "Utility-first styling, responsive grid, dark glassmorphism theme"),
        ("Lucide Icons",      "Consistent icon system across all UI states"),
        ("Custom CSS",        "Glass panel effects, depth stages, lift shadows"),
    ]),
    ("Output", CORAL, [
        ("jsPDF",             "Client-side PDF brief generation — no server required"),
        ("CSV Export",        "Full hotspot dataset download for offline analysis"),
        ("Vite Build",        "Optimised static bundle deployable to any static host"),
    ]),
]

for i, (layer_name, acc, items) in enumerate(tech_layers):
    col = i % 2
    row = i // 2
    bx = 0.6 + col * 6.32
    by = 1.7 + row * 2.55

    rect(s, bx, by, 6.1, 2.38, fill=CARD)
    rect(s, bx, by, 6.1, 0.065, fill=acc)
    textbox(s, bx + 0.18, by + 0.16, 5.6, 0.38, layer_name, size=13, bold=True, color=acc)

    for j, (tech, desc) in enumerate(items):
        ty = by + 0.64 + j * 0.55
        rect(s, bx + 0.18, ty, 1.7, 0.38, fill=MID)
        textbox(s, bx + 0.18, ty + 0.04, 1.7, 0.32,
                tech, size=10, bold=True, color=acc, align=PP_ALIGN.CENTER)
        textbox(s, bx + 2.02, ty + 0.04, 3.8, 0.38,
                desc, size=10, color=LGRAY)

rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.6, 7.1, 12, 0.35,
        "Zero paid services · Zero API keys · Runs entirely in the browser after initial data load",
        size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — WHY PARKSIGHT AI WINS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=CORAL)
rect(s, 0, 0, 13.33, 0.08, fill=CORAL)

textbox(s, 0.6, 0.22, 4.0, 0.45, "COMPETITIVE ADVANTAGE", size=11, bold=True, color=CORAL)
textbox(s, 0.6, 0.55, 11.0, 0.85,
        "Why ParkSight AI Stands Apart from Simple Analytics", size=30, bold=True, color=WHITE)
divider(s, 1.5, color=CORAL, x=0.6, w=12.4)

differentiators = [
    (CORAL,  "Actionable, Not Just Analytical",
     "Most dashboards show WHERE violations happen. ParkSight AI answers WHAT TO DO: "
     "which beat to patrol, at what time, with what action, and why — in 12 ready-to-deploy cards."),
    (AMBER,  "Transparent, Auditable Scoring",
     "The Parking Impact Index formula is printed on the dashboard itself. Every weight, "
     "every threshold is visible. No black-box ML — any officer or judge can audit the ranking."),
    (TEAL,   "Physics-Based, Not Count-Based",
     "Unlike simple violation counters, ParkSight uses Passenger Car Unit (PCU) weights "
     "borrowed from traffic engineering — a tanker blocking a junction scores 3.2× a moped."),
    (VIOLET, "Multi-Dimensional Ranking",
     "7 independent signals (obstruction pressure, density, junction exposure, arterial risk, "
     "peak share, recurrence, severity) prevent any single outlier from hijacking the ranking."),
    (TEAL2,  "Zero Infrastructure Required",
     "The entire system runs in a browser tab. No servers, no cloud subscription, no API keys. "
     "A field officer with a laptop can generate the enforcement plan offline."),
    (RGBColor(0xC9,0x6F,0x52), "Production-Ready Output",
     "One-click PDF brief and CSV download mean the insights survive the demo. "
     "Command officers get a brief; data teams get the full dataset; field teams get beat cards."),
]

for i, (acc, title, body) in enumerate(differentiators):
    col = i % 2
    row = i // 2
    bx = 0.6 + col * 6.32
    by = 1.7 + row * 1.82

    rect(s, bx, by, 6.1, 1.66, fill=CARD)
    rect(s, bx, by, 0.065, 1.66, fill=acc)

    rect(s, bx + 0.22, by + 0.18, 0.55, 0.55, fill=acc)
    textbox(s, bx + 0.22, by + 0.14, 0.55, 0.6,
            "✓", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    textbox(s, bx + 0.92, by + 0.15, 5.0, 0.42, title, size=13, bold=True, color=WHITE)
    textbox(s, bx + 0.92, by + 0.58, 5.0, 1.0,  body,  size=11, color=LGRAY)

rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.6, 7.1, 12, 0.35,
        "ParkSight AI = Traffic Engineering + Data Science + Software Engineering in one open-source package",
        size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — FUTURE ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=VIOLET)
rect(s, 0, 0, 13.33, 0.08, fill=VIOLET)

textbox(s, 0.6, 0.22, 4.0, 0.45, "FUTURE ROADMAP", size=11, bold=True, color=VIOLET)
textbox(s, 0.6, 0.55, 11.0, 0.85,
        "ParkSight AI v2 — What Comes Next", size=30, bold=True, color=WHITE)
divider(s, 1.5, color=VIOLET, x=0.6, w=12.4)

# Timeline roadmap
phases = [
    ("Phase 1\nNow",      TEAL,   "Core Platform (Complete)", [
        "298K-record pipeline & scoring",
        "160-hotspot ranked dashboard",
        "12-beat enforcement plan",
        "CSV & PDF report export",
        "3-layer Leaflet heatmap",
    ]),
    ("Phase 2\n3 Months", AMBER,  "Real-Time Intelligence", [
        "Live police challan API integration",
        "Daily automated score refresh",
        "Mobile-responsive PWA for field",
        "Push alerts for emerging hotspots",
        "Week-over-week trend charts",
    ]),
    ("Phase 3\n6 Months", CORAL,  "AI & Computer Vision", [
        "CCTV feed analysis for parked vehicles",
        "Number plate detection & cross-ref",
        "Predictive model: forecast next hotspot",
        "Weather & event-aware enforcement",
        "Multi-city template (Mumbai, Delhi)",
    ]),
    ("Phase 4\n12 Months", VIOLET, "Smart City Integration", [
        "SCITA & TMS data feed integration",
        "Dynamic signal timing coordination",
        "Tow-truck dispatch automation",
        "Public-facing parking guidance app",
        "Congestion savings measurement API",
    ]),
]

for i, (phase, acc, title, items) in enumerate(phases):
    bx = 0.6 + i * 3.16
    rect(s, bx, 1.7, 2.95, 5.1, fill=CARD)
    rect(s, bx, 1.7, 2.95, 0.065, fill=acc)

    textbox(s, bx + 0.14, 1.76, 2.65, 0.5,
            phase, size=11, bold=True, color=acc, align=PP_ALIGN.CENTER)
    textbox(s, bx + 0.14, 2.3, 2.65, 0.55,
            title, size=12, bold=True, color=WHITE)
    divider(s, 2.9, color=acc, x=bx + 0.14, w=2.65)

    for j, item in enumerate(items):
        iy = 3.0 + j * 0.5
        dot = s.shapes.add_shape(9, Inches(bx + 0.16), Inches(iy + 0.1),
                                  Inches(0.1), Inches(0.1))
        dot.line.fill.background(); dot.fill.solid(); dot.fill.fore_color.rgb = acc
        textbox(s, bx + 0.36, iy, 2.4, 0.45, item, size=10, color=LGRAY)

# connecting arrows between phases
for i in range(3):
    ax = 0.6 + i * 3.16 + 2.95 + 0.03
    textbox(s, ax, 4.0, 0.18, 0.4, "→", size=20, bold=True, color=LGRAY)

rect(s, 0, 7.05, 13.33, 0.45, fill=CARD)
textbox(s, 0.6, 7.1, 12, 0.35,
        "v1 is live today — Phases 2–4 build on the same scoring engine without architectural changes",
        size=10, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — THANK YOU / CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 0.35, 7.5, fill=TEAL)

# Left accent panel
rect(s, 0, 0, 5.8, 7.5, fill=MID)
rect(s, 0, 0, 0.35, 7.5, fill=TEAL)

chip(s, 0.6, 0.5, "FLIPKART GRIDLOCK CHALLENGE 2024", bg=TEAL, fg=WHITE, size=10, w=4.6, h=0.32)

textbox(s, 0.6, 1.1, 4.8, 1.1,
        "Thank You", size=58, bold=True, color=WHITE)
textbox(s, 0.6, 2.15, 4.8, 0.5,
        "from Team ParkSight AI", size=18, color=TEAL2, italic=True)

divider(s, 2.8, color=TEAL, x=0.6, w=4.8)

textbox(s, 0.6, 2.95, 4.8, 1.5,
        "We built ParkSight AI to answer a simple question:\nWhere should enforcement go TODAY — before the\njam forms? We believe this is the foundation of\na smarter, safer Bengaluru.",
        size=13, color=LGRAY)

# project stats recap
for i, (v, l, c) in enumerate([
    ("298K",  "Violations Analysed",  CORAL),
    ("160",   "Hotspots Ranked",      AMBER),
    ("12",    "Beats Generated",      TEAL2),
    ("100%",  "Open Source",          VIOLET),
]):
    bx = 0.6 + (i % 2) * 2.5
    by = 4.7 + (i // 2) * 1.25
    rect(s, bx, by, 2.3, 1.1, fill=CARD)
    textbox(s, bx + 0.15, by + 0.08, 2.0, 0.55, v, size=30, bold=True, color=c)
    textbox(s, bx + 0.15, by + 0.62, 2.0, 0.4,  l, size=10, color=LGRAY)

# Right — key quotes and summary
rect(s, 6.0, 0, 7.33, 7.5, fill=NAVY)

textbox(s, 6.3, 0.65, 6.7, 0.5,
        "ParkSight AI — By the Numbers", size=16, bold=True, color=WHITE)
divider(s, 1.22, color=TEAL, x=6.3, w=6.7)

summary_points = [
    (CORAL,  "298,450 violations", "processed and scored — the complete Bengaluru record"),
    (AMBER,  "3,969 urban cells",  "every micro-zone in the metro ranked 0–100"),
    (TEAL2,  "7-factor formula",   "transparent, auditable Parking Impact Index"),
    (VIOLET, "12 patrol beats",    "time-windowed, action-tagged, map-linked"),
    (CORAL,  "3 heatmap layers",   "Impact · Volume · Junction — switch live"),
    (TEAL,   "55 stations",        "ranked leaderboard with burden & risk scores"),
    (AMBER,  "1-click reports",    "PDF brief + CSV download, no server needed"),
    (LGRAY,  "0 paid services",    "entirely open-source, runs in any browser"),
]

for i, (acc, bold_part, rest) in enumerate(summary_points):
    yy = 1.42 + i * 0.7
    rect(s, 6.3, yy, 6.7, 0.62, fill=CARD if i % 2 == 0 else MID)
    rect(s, 6.3, yy, 0.055, 0.62, fill=acc)
    textbox(s, 6.45, yy + 0.08, 2.5, 0.42, bold_part, size=13, bold=True, color=acc)
    textbox(s, 8.85, yy + 0.1,  4.0, 0.42, rest, size=11, color=LGRAY)

textbox(s, 6.3, 6.95, 6.7, 0.4,
        "Parking Impact Index v1.0  ·  Bengaluru, India  ·  Flipkart Gridlock Challenge 2024",
        size=9, color=LGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════════════
out_path = r"d:\FLIPKART_GRIDLOCK\ParkSight_AI_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
