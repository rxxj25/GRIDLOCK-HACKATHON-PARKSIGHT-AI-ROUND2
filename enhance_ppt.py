"""
Enhances ParkSight_AI_Pitch_Deck.pptx — improves slides 2, 6, 9, 10.
Adds visual panels, stat cards, accent bars, and formula boxes on top of
(or behind) the existing shapes without touching slides that are already strong.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC  = r'd:\FLIPKART_GRIDLOCK\outputs\ParkSight_AI_Pitch_Deck.pptx'
DEST = r'd:\FLIPKART_GRIDLOCK\outputs\ParkSight_AI_Pitch_Deck_v2.pptx'

prs = Presentation(SRC)

# ── colour palette (matches the app + existing deck) ──────────────────────────
NAVY   = RGBColor(0x17, 0x20, 0x2A)
TEAL   = RGBColor(0x04, 0x75, 0x6F)
TEAL2  = RGBColor(0x05, 0x96, 0x8E)
CORAL  = RGBColor(0xD9, 0x3D, 0x4A)
AMBER  = RGBColor(0xE0, 0x9B, 0x2D)
VIOLET = RGBColor(0x61, 0x57, 0xA8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0x65, 0x70, 0x7D)
DCARD  = RGBColor(0x1A, 0x2B, 0x3C)   # dark card fill
CLIGHT = RGBColor(0xF4, 0xF8, 0xFC)   # very light card
C_GRN  = RGBColor(0xDF, 0xF4, 0xE8)   # light green text
CORAL_L= RGBColor(0xFD, 0xF0, 0xF1)   # light coral tint
TEAL_L = RGBColor(0xF0, 0xFA, 0xF9)   # light teal tint
DARK_T = RGBColor(0x0B, 0x3D, 0x38)   # dark teal for result bar


# ── helpers ───────────────────────────────────────────────────────────────────

def send_to_back(slide):
    """Move the last-added shape behind all pre-existing shapes."""
    sp_tree = slide.shapes._spTree
    last = sp_tree[-1]
    sp_tree.remove(last)
    sp_tree.insert(2, last)


def R(slide, x, y, w, h, fill, back=False):
    """Add a solid rectangle."""
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if back:
        send_to_back(slide)
    return shp


def T(slide, x, y, w, h, text, size=12, bold=False, color=NAVY,
      align=PP_ALIGN.LEFT, italic=False):
    """Add a text box."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb


def stat_card(slide, x, y, value, label, accent, w=2.68, h=1.3):
    """Dark stat card with left accent bar."""
    R(slide, x, y, w, h, DCARD)
    R(slide, x, y, 0.07, h, accent)
    T(slide, x + 0.2, y + 0.08, w - 0.3, 0.54, value, size=26, bold=True, color=WHITE)
    T(slide, x + 0.2, y + 0.65, w - 0.3, 0.55, label, size=9.5,
      color=RGBColor(0xC4, 0xCF, 0xD8))


def card_with_stripe(slide, x, y, w, h, accent, back=True):
    """Light card with a coloured top stripe — send to back so existing text floats above."""
    R(slide, x, y, w, h, CLIGHT, back=back)
    R(slide, x, y, w, 0.055, accent)


def bullet_chip(slide, x, y, text, accent):
    """Small coloured dot + bold label for a bullet row."""
    R(slide, x, y + 0.07, 0.11, 0.11, accent)
    T(slide, x + 0.18, y, 2.3, 0.3, text, size=10, bold=True, color=accent)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — THE OPERATIONAL GAP
#  Enhancement: replace sparse right-side "Today" block with a dark stat panel;
#               add light card + teal accent behind the left bullets.
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides[1]

# ── left side: light card behind bullets (sent to back → text stays visible) ──
R(s2, 0.76, 2.06, 6.25, 2.65, CLIGHT, back=True)
# Teal accent bar on far left (non-overlapping with text at x=1.17)
R(s2, 0.76, 2.06, 0.07, 2.65, TEAL)

# ── right side: dark panel covers sparse "Today" text (stays in front) ────────
R(s2, 7.12, 1.92, 5.98, 5.0, NAVY)
R(s2, 7.12, 1.92, 5.98, 0.06, TEAL)  # teal top strip

# Section label
T(s2, 7.32, 2.02, 5.5, 0.28,
  "THE SCALE OF THE PROBLEM", size=9, bold=True, color=TEAL2)

# 2 × 2 stat cards
stat_card(s2, 7.32, 2.4,  "Rs.3,700 Cr", "Annual traffic loss,\nBengaluru city",    CORAL)
stat_card(s2, 10.18, 2.4, "298,450",      "Violations recorded\nin just 5 months",  AMBER)
stat_card(s2, 7.32, 3.88, "50.5%",        "Cases at junctions\n& crossings",        TEAL)
stat_card(s2, 10.18, 3.88,"39.6%",        "During peak movement\nhours (8-11 AM, 5-9 PM)", VIOLET)

# Result bar at the bottom of the dark panel
R(s2, 7.12, 5.62, 5.98, 0.72, DARK_T)
T(s2, 7.3, 5.7, 5.65, 0.56,
  "Result: Enforcement effort and congestion pain don't align. ParkSight AI fixes this.",
  size=11, bold=True, color=C_GRN)

# Thin divider between left and right sections
R(s2, 7.04, 2.06, 0.04, 3.58, TEAL)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — EXPLAINABLE AI / PARKING IMPACT INDEX
#  Enhancement: coloured top-stripe cards behind each pipeline step;
#               prominent dark formula box; weight visualisation bar row.
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides[5]

# 4 step cards — sent to back so existing step text floats on top
step_x      = [0.78, 3.78, 6.78, 9.78]
step_colors = [TEAL, AMBER, CORAL, VIOLET]
for cx, col in zip(step_x, step_colors):
    card_with_stripe(s6, cx, 2.1, 2.75, 2.05, col, back=True)

# Prominent formula panel (covers existing plain formula text area ~y=4.98)
R(s6, 0.75, 4.68, 12.08, 0.92, NAVY)
R(s6, 0.75, 4.68, 0.07,  0.92, AMBER)
T(s6, 0.93, 4.72, 1.55, 0.3, "Formula", size=9.5, bold=True, color=AMBER)
T(s6, 2.42, 4.72, 9.9, 0.62,
  "Score = 100 x (0.34 obstruction + 0.18 density + 0.15 junction + "
  "0.13 arterial + 0.10 peak + 0.06 recurrence + 0.04 severity)",
  size=11, bold=True, color=WHITE)

# Weight visualisation strip (covers existing explanation text ~y=6.02)
R(s6, 0.75, 5.72, 12.08, 1.08, CLIGHT)
R(s6, 0.75, 5.72, 0.07,  1.08, TEAL)
T(s6, 0.95, 5.76, 3.5, 0.28,
  "Why explainable?  No speed/congestion labels in data — every weight is auditable.",
  size=9.5, italic=True, color=LGRAY)

weights = [
    ("Obstruction", 34, CORAL),
    ("Density",     18, AMBER),
    ("Junction",    15, TEAL),
    ("Arterial",    13, RGBColor(0x56, 0xA3, 0xA6)),
    ("Peak",        10, VIOLET),
    ("Recurr.",      6, LGRAY),
    ("Severity",     4, RGBColor(0x88, 0x90, 0x63)),
]
MAX_BAR = 1.58   # bar width at 34%
bx = 0.88
for wname, wpct, wcol in weights:
    T(s6, bx, 6.08, 1.6, 0.26, f"{wname}  {wpct}%", size=8.5, bold=True, color=NAVY)
    R(s6, bx, 6.38, MAX_BAR * wpct / 34, 0.2, wcol)
    bx += 1.68


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — SYSTEM ARCHITECTURE
#  Enhancement: coloured node cards with accent stripes;
#               arrow indicators between nodes; enhanced "hosted" strip.
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides[8]

# 4 node cards — sent to back so existing node text floats above
node_x      = [0.75, 3.84, 6.88, 9.94]
node_colors = [AMBER, TEAL, CORAL, VIOLET]
node_tech   = ["CSV  |  298K rows", "Python 3.11  ETL", "React + Leaflet", "12 Beats  |  CSV  |  PDF"]

for i, (nx, nc, nt) in enumerate(zip(node_x, node_colors, node_tech)):
    card_with_stripe(s9, nx, 2.2, 2.58, 2.1, nc, back=True)
    # Small tech chip at bottom of card (in front of card background)
    R(s9, nx, 4.1, 2.58, 0.22, nc)
    T(s9, nx + 0.08, 4.11, 2.42, 0.2, nt,
      size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow indicators between nodes (in the ~0.2" gaps)
for i in range(3):
    ax = node_x[i] + 2.61
    T(s9, ax, 3.05, 0.22, 0.38, ">", size=20, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)

# Enhanced "Hosted online" panel (existing text at y=5.21–5.52 floats on top)
R(s9, 0.75, 5.08, 12.25, 0.72, CLIGHT)
R(s9, 0.75, 5.08, 0.07,  0.72, TEAL)

# Step counter chips below header
for i, (label, col) in enumerate([
    ("Step 1: Ingest", AMBER), ("Step 2: Score", TEAL),
    ("Step 3: Visualise", CORAL), ("Step 4: Output", VIOLET),
]):
    R(s9, 0.85 + i * 3.14, 1.82, 2.88, 0.27, col)
    T(s9, 0.85 + i * 3.14, 1.84, 2.88, 0.23,
      label, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — IMPACT (BEFORE / AFTER)
#  Enhancement: coral-tinted "Before" panel; teal-tinted "After" panel;
#               dramatic VS divider; dark roadmap strip.
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides[9]

# Before panel (light coral) — sent to back so existing dark text stays visible
R(s10, 0.75, 1.92, 5.82, 3.72, CORAL_L, back=True)
R(s10, 0.75, 1.92, 0.07, 3.72, CORAL,   back=False)  # accent bar (in front)

# After panel (light teal) — sent to back
R(s10, 7.25, 1.92, 5.82, 3.72, TEAL_L, back=True)
R(s10, 7.25, 1.92, 0.07, 3.72, TEAL,   back=False)   # accent bar (in front)

# VS divider pill centred in gap (x=6.55–7.25)
R(s10, 6.52, 2.7, 0.58, 1.7, NAVY)
T(s10, 6.52, 3.35, 0.58, 0.5,
  "VS", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Column header chips
R(s10, 0.82, 1.95, 1.88, 0.28, CORAL)
T(s10, 0.82, 1.96, 1.88, 0.26,
  "BEFORE", size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

R(s10, 7.32, 1.95, 2.22, 0.28, TEAL)
T(s10, 7.32, 1.96, 2.22, 0.26,
  "AFTER  PARKSIGHT AI", size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Check-mark row (visual upgrade on the "After" side)
after_items = [
    "Impact-ranked hotspots",
    "Station-wise navigation",
    "12-beat enforcement plan",
    "Downloadable PDF/CSV brief",
]
for i, item in enumerate(after_items):
    yy = 2.88 + i * 0.52
    R(s10, 7.72, yy + 0.05, 0.24, 0.24, TEAL)
    T(s10, 7.72, yy + 0.03, 0.24, 0.24,
      "ok", size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Roadmap strip (covers existing footer-level roadmap text ~y=6.15)
R(s10, 0.75, 5.82, 12.25, 0.7, NAVY)
R(s10, 0.75, 5.82, 0.07,  0.7, AMBER)
T(s10, 0.95, 5.9, 11.85, 0.55,
  "Roadmap: integrate live CCTV/YOLO illegal-parking detection and "
  "traffic-speed labels for supervised congestion prediction.",
  size=11, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  GLOBAL: thin teal accent stripe on slides 3, 4, 7, 8
#          (adds visual consistency with enhanced slides)
# ══════════════════════════════════════════════════════════════════════════════
for idx in [2, 3, 6, 7]:   # slides 3, 4, 7, 8 (0-indexed)
    slide = prs.slides[idx]
    R(slide, 0.0, 0.0, 13.33, 0.045, TEAL)   # very thin top bar

prs.save(DEST)
print(f"Saved: {DEST}")
print(f"Total slides: {len(prs.slides)}")
