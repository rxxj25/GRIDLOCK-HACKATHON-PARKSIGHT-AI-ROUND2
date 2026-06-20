import sys
from pptx import Presentation

def main():
    deck_path = "outputs/ParkSight_AI_Pitch_Deck.pptx"
    print(f"Loading presentation from {deck_path}...")
    prs = Presentation(deck_path)

    # 1. Update Slide 6 index formula
    # Slide 6 is at index 5
    slide6 = prs.slides[5]
    formula_old = "0.34 obstruction + 0.18 density + 0.15 junction + 0.13 arterial + 0.10 peak + recurrence + severity"
    formula_new = "0.34 obstruction + 0.18 density + 0.15 junction + 0.13 arterial + 0.10 peak + 0.06 recurrence + 0.04 severity"
    
    found_formula = False
    for shape in slide6.shapes:
        if shape.has_text_frame and formula_old in shape.text:
            print("Found formula on Slide 6. Updating...")
            for paragraph in shape.text_frame.paragraphs:
                if formula_old in paragraph.text:
                    for run in paragraph.runs:
                        if formula_old in run.text:
                            run.text = run.text.replace(formula_old, formula_new)
                            found_formula = True
                    if not found_formula:
                        paragraph.text = paragraph.text.replace(formula_old, formula_new)
                        found_formula = True
            break
            
    if not found_formula:
        # Try a substring search in case of slight whitespace mismatches
        for shape in slide6.shapes:
            if shape.has_text_frame and "0.34 obstruction" in shape.text and "recurrence + severity" in shape.text:
                print("Found formula with slight mismatch. Updating text...")
                for paragraph in shape.text_frame.paragraphs:
                    if "0.34 obstruction" in paragraph.text:
                        paragraph.text = "0.34 obstruction + 0.18 density + 0.15 junction + 0.13 arterial + 0.10 peak + 0.06 recurrence + 0.04 severity"
                        found_formula = True
                        break

    # 2. Update Slide 8 beat names to differentiate them
    # Slide 8 is at index 7
    slide8 = prs.slides[7]
    
    # We want to identify the shapes for Beat 1, 2, 3, 4 based on their vertical positions (y-coordinate)
    # or just list all text shapes containing the targets and sort them by vertical position.
    targets = ["BTP051 - Safina Plaza Junction", "BTP040 - Elite Junction"]
    text_shapes = []
    
    for shape in slide8.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text in targets:
                text_shapes.append(shape)
                
    # Sort from top to bottom based on shape top coordinate
    text_shapes.sort(key=lambda s: s.top)
    
    print(f"Found {len(text_shapes)} target beat names on Slide 8 to update:")
    for i, shape in enumerate(text_shapes):
        print(f"  Shape {i+1} at top={shape.top}: '{shape.text}'")
        
    # We expect 4 shapes sorted top-to-bottom:
    # 1. Safina Plaza Junction (Beat 1) -> Cell A
    # 2. Elite Junction (Beat 2) -> Cell A
    # 3. Elite Junction (Beat 3) -> Cell B
    # 4. Safina Plaza Junction (Beat 4) -> Cell B
    
    if len(text_shapes) >= 4:
        # Shape 1: Beat 1
        p = text_shapes[0].text_frame.paragraphs[0]
        p.text = "BTP051 - Safina Plaza (Cell A)"
        
        # Shape 2: Beat 2
        p = text_shapes[1].text_frame.paragraphs[0]
        p.text = "BTP040 - Elite Junction (Cell A)"
        
        # Shape 3: Beat 3
        p = text_shapes[2].text_frame.paragraphs[0]
        p.text = "BTP040 - Elite Junction (Cell B)"
        
        # Shape 4: Beat 4
        p = text_shapes[3].text_frame.paragraphs[0]
        p.text = "BTP051 - Safina Plaza (Cell B)"
        
        print("Updated Slide 8 beat labels to indicate distinct grid cells.")
    else:
        print("Warning: Could not find all 4 beat name shapes to update on Slide 8.")

    prs.save(deck_path)
    print(f"Successfully saved changes to {deck_path}!")

if __name__ == "__main__":
    main()
